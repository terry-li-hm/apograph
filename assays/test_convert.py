"""Basic conversion test — extract + convert round trip."""

import tempfile
from pathlib import Path

from apograph.extract import extract_from_html
from apograph.convert import convert


SAMPLE_HTML = """<!DOCTYPE html><html><body>
<div class="slide" style="width:800px;height:450px;background:#fafafa;position:relative;font-family:sans-serif">
  <h1 style="font-size:24px;font-weight:700;color:#333;margin:20px">Test Title</h1>
  <p style="font-size:12px;color:#777;margin:0 20px">Subtitle text here</p>
  <div style="background:#ddd;height:1px;width:760px;margin:10px 20px"></div>
  <div style="width:40px;height:40px;border-radius:50%;background:#eee;margin:10px 20px"></div>
</div>
</body></html>"""


def test_basic_conversion():
    with tempfile.NamedTemporaryFile(suffix=".html", mode="w", delete=False) as tmp:
        tmp.write(SAMPLE_HTML)
        html_path = Path(tmp.name)

    out_path = html_path.with_suffix(".pptx")
    try:
        data = extract_from_html(html_path)
        assert data.width_px > 0
        assert data.height_px > 0
        assert len(data.elements) > 0

        result = convert(data, out_path)
        assert result.exists()
        assert result.stat().st_size > 1000  # non-trivial PPTX

        from pptx import Presentation
        prs = Presentation(str(result))
        slide = prs.slides[0]
        assert len(slide.shapes) >= 3  # title + subtitle + at least one shape
    finally:
        html_path.unlink(missing_ok=True)
        out_path.unlink(missing_ok=True)


def test_hybrid_conversion():
    with tempfile.NamedTemporaryFile(suffix=".html", mode="w", delete=False) as tmp:
        tmp.write(SAMPLE_HTML)
        html_path = Path(tmp.name)

    out_path = html_path.with_suffix(".pptx")
    try:
        data = extract_from_html(html_path, hybrid=True)
        assert data.background_image is not None
        assert len(data.background_image) > 100  # non-trivial PNG

        result = convert(data, out_path, hybrid=True)
        assert result.exists()

        from pptx import Presentation
        prs = Presentation(str(result))
        slide = prs.slides[0]
        # Hybrid: background image + text shapes (no native rules/circles)
        has_picture = any(hasattr(s, "image") for s in slide.shapes)
        assert has_picture, "Hybrid mode should have background image"
    finally:
        html_path.unlink(missing_ok=True)
        out_path.unlink(missing_ok=True)
