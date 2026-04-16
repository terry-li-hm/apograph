"""Convert extracted HTML layout to PPTX."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

from apograph.extract import SlideData

# Default PPTX slide dimensions (16:9)
DEFAULT_SLIDE_W_IN = 13.333
DEFAULT_SLIDE_H_IN = 7.5


def _px_to_emu(val_px: float, scale: float) -> int:
    """Convert CSS pixels to EMU at the given scale (inches per pixel)."""
    return int(val_px * scale * 914400)


def _parse_rgb(css_color: str) -> RGBColor:
    """Parse 'rgb(r, g, b)' or 'rgba(r, g, b, a)' to RGBColor."""
    css_color = css_color.strip()
    if css_color.startswith("rgba"):
        parts = css_color[5:-1].split(",")
    elif css_color.startswith("rgb"):
        parts = css_color[4:-1].split(",")
    else:
        return RGBColor(0x33, 0x33, 0x33)
    red, green, blue = int(parts[0].strip()), int(parts[1].strip()), int(parts[2].strip())
    return RGBColor(red, green, blue)


def convert(
    data: SlideData,
    output_path: Path,
    *,
    slide_width_in: float = DEFAULT_SLIDE_W_IN,
    slide_height_in: float = DEFAULT_SLIDE_H_IN,
    image_base_dir: Path | None = None,
    accent_color: str | None = None,
    accent_height_px: float = 3.0,
    font_override: str | None = None,
) -> Path:
    """Build a PPTX from extracted HTML slide data.

    Args:
        data: Extracted layout from HTML.
        output_path: Where to save the .pptx file.
        slide_width_in: PPTX slide width in inches.
        slide_height_in: PPTX slide height in inches.
        image_base_dir: Directory to resolve relative image paths.
        accent_color: CSS color for top accent bar (extracted from ::before).
        accent_height_px: Height of accent bar in CSS pixels.
        font_override: Force all text to this font family.
    """
    scale = slide_width_in / data.width_px  # inches per CSS pixel

    prs = Presentation()
    prs.slide_width = Emu(int(slide_width_in * 914400))
    prs.slide_height = Emu(int(slide_height_in * 914400))
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background — detect from first element or default
    bg_color = RGBColor(0xFC, 0xFB, 0xFA)
    for elem in data.elements:
        if elem.get("depth") == 0:
            bg_css = elem.get("backgroundColor", "")
            if bg_css and "rgba(0, 0, 0, 0)" not in bg_css:
                bg_color = _parse_rgb(bg_css)
            break
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    # Accent bar (::before pseudo-element not captured by JS)
    if accent_color:
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, _px_to_emu(accent_height_px, scale))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _parse_rgb(accent_color)
        accent.line.fill.background()

    for elem in data.elements:
        left = _px_to_emu(elem["x"], scale)
        top = _px_to_emu(elem["y"], scale)
        width = _px_to_emu(elem["w"], scale)
        height = _px_to_emu(elem["h"], scale)
        cls = elem.get("cls", "")
        tag = elem.get("tag", "")
        text = elem.get("text")

        # --- Images ---
        if elem.get("isImage"):
            src = elem.get("src", "")
            if src and image_base_dir:
                img_path = (image_base_dir / Path(src).name)
                if img_path.exists():
                    slide.shapes.add_picture(str(img_path), left, top, width, height)
            continue

        # --- Placeholder circles ---
        if elem.get("isPlaceholder"):
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xED, 0xED, 0xED)
            shape.line.fill.background()
            continue

        # --- Horizontal rules ---
        if "rule" in cls and tag == "div":
            # Enforce minimum 0.75pt height so the line is visible in PPTX
            min_height = Pt(0.75)
            rule_height = max(height, min_height)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, rule_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xD7, 0xD8, 0xD6)
            shape.line.fill.background()
            continue

        # --- Text ---
        if text and len(text.strip()) > 0:
            box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = text_frame.margin_right = 0
            text_frame.margin_top = text_frame.margin_bottom = 0
            para = text_frame.paragraphs[0]

            display_text = text.strip()
            if elem.get("textTransform") == "uppercase":
                display_text = display_text.upper()

            # Handle line breaks (<br> → multiple paragraphs)
            lines = display_text.split("\n")
            para.text = lines[0]

            # Font sizing: CSS px → PPTX pt at slide scale
            font_size_px = elem.get("fontSize", 12)
            pt_size = font_size_px * scale * 72
            para.font.size = Pt(pt_size)

            para.font.color.rgb = _parse_rgb(elem.get("color", "rgb(51,51,51)"))

            weight = str(elem.get("fontWeight", "400"))
            para.font.bold = weight in ("700", "bold", "800", "900")
            para.font.italic = elem.get("fontStyle") == "italic"
            para.font.name = font_override or "Century Gothic"

            # Add remaining lines as new paragraphs
            for extra_line in lines[1:]:
                new_para = text_frame.add_paragraph()
                new_para.text = extra_line.strip()
                new_para.font.size = Pt(pt_size)
                new_para.font.color.rgb = _parse_rgb(elem.get("color", "rgb(51,51,51)"))
                new_para.font.bold = weight in ("700", "bold", "800", "900")
                new_para.font.italic = elem.get("fontStyle") == "italic"
                new_para.font.name = font_override or "Century Gothic"
                new_para.space_after = Pt(0)
                new_para.space_before = Pt(0)

            # Line-height correction: PPT renders ~20% tighter than CSS
            css_line_height = elem.get("lineHeight", "")
            if css_line_height and css_line_height not in ("normal", ""):
                try:
                    lh_px = float(css_line_height.replace("px", ""))
                    lh_ratio = lh_px / font_size_px if font_size_px else 1.0
                    # PPT line_spacing as float = multiple of single-space
                    para.line_spacing = lh_ratio * 0.85
                except (ValueError, ZeroDivisionError):
                    pass

            # Letter-spacing: CSS px → PPT charSpacing (hundredths of a point)
            # Applied to the default run properties on the paragraph
            css_letter_spacing = elem.get("letterSpacing", "")
            if css_letter_spacing and css_letter_spacing not in ("normal", "0px"):
                try:
                    ls_px = float(css_letter_spacing.replace("px", ""))
                    if ls_px != 0:
                        char_spacing = int(ls_px * 72 / 96 * 100)
                        from lxml import etree
                        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                        # Set on default run properties
                        def_rpr = para._p.find("a:pPr/a:defRPr", nsmap)
                        if def_rpr is None:
                            p_pr = para._p.find("a:pPr", nsmap)
                            if p_pr is None:
                                p_pr = etree.SubElement(para._p, f"{{{nsmap['a']}}}pPr")
                            def_rpr = etree.SubElement(p_pr, f"{{{nsmap['a']}}}defRPr")
                        def_rpr.set("spc", str(char_spacing))
                        # Also set on existing runs
                        for run_elem in para._p.findall("a:r/a:rPr", nsmap):
                            run_elem.set("spc", str(char_spacing))
                except ValueError:
                    pass

            para.space_after = Pt(0)
            para.space_before = Pt(0)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
